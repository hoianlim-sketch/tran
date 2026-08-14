"""DOCX·PPTX 문단을 제자리 번역해 서식·표·슬라이드를 유지한다."""

from __future__ import annotations

import io
from collections.abc import Callable
from pathlib import Path

from docx import Document
from docx.text.paragraph import Paragraph

from src.extract_errors import ExtractError

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:  # pragma: no cover
    Presentation = None
    MSO_SHAPE_TYPE = None

TranslateUnits = Callable[[list[str]], list[str]]


def can_preserve(suffix: str) -> bool:
    return suffix.lower() in {".docx", ".pptx"}


def translate_preserving(
    filename: str,
    data: bytes,
    translate_units: TranslateUnits,
) -> bytes:
    suffix = Path(filename).suffix.lower()
    if suffix == ".docx":
        return translate_docx(data, translate_units)
    if suffix == ".pptx":
        return translate_pptx(data, translate_units)
    raise ExtractError(f"서식 유지 번역을 지원하지 않는 형식입니다: {suffix}")


def translate_docx(data: bytes, translate_units: TranslateUnits) -> bytes:
    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractError(f"DOCX를 읽지 못했습니다: {exc}") from exc

    paragraphs = list(_docx_paragraphs(document))
    originals = [p.text for p in paragraphs]
    translated = translate_units(originals)
    if len(translated) != len(paragraphs):
        raise ExtractError("서식 유지 번역 결과가 문단 수와 맞지 않습니다.")
    for paragraph, new_text, old_text in zip(paragraphs, translated, originals):
        if old_text.strip() and new_text.strip() and new_text != old_text:
            _set_paragraph_text(paragraph, new_text)

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def translate_pptx(data: bytes, translate_units: TranslateUnits) -> bytes:
    if Presentation is None:
        raise ExtractError("PPTX 지원을 위해 python-pptx가 필요합니다.")
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractError(f"PPTX를 읽지 못했습니다: {exc}") from exc

    targets: list[tuple[object, str]] = []
    for slide in presentation.slides:
        _collect_pptx_shapes(slide.shapes, targets)

    originals = [item[1] for item in targets]
    translated = translate_units(originals)
    if len(translated) != len(targets):
        raise ExtractError("서식 유지 번역 결과가 텍스트 상자 수와 맞지 않습니다.")
    for (paragraph, old_text), new_text in zip(targets, translated):
        if old_text.strip() and new_text.strip() and new_text != old_text:
            _set_pptx_paragraph(paragraph, new_text)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _docx_paragraphs(document: Document) -> list[Paragraph]:
    items: list[Paragraph] = []
    items.extend(document.paragraphs)
    for table in document.tables:
        items.extend(_table_paragraphs(table))
    for section in document.sections:
        items.extend(section.header.paragraphs)
        items.extend(section.footer.paragraphs)
    return items


def _table_paragraphs(table) -> list[Paragraph]:
    items: list[Paragraph] = []
    for row in table.rows:
        for cell in row.cells:
            items.extend(cell.paragraphs)
            for nested in cell.tables:
                items.extend(_table_paragraphs(nested))
    return items


def _set_paragraph_text(paragraph: Paragraph, text: str) -> None:
    if not paragraph.runs:
        paragraph.add_run(text)
        return
    paragraph.runs[0].text = text
    for run in paragraph.runs[1:]:
        run.text = ""


def _collect_pptx_shapes(shapes, targets: list[tuple[object, str]]) -> None:
    for shape in shapes:
        if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_pptx_shapes(shape.shapes, targets)
            continue
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs) or paragraph.text
                targets.append((paragraph, text))
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs) or paragraph.text
                        targets.append((paragraph, text))


def _set_pptx_paragraph(paragraph, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
        return
    paragraph.text = text
