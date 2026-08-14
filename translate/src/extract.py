"""PDF, DOCX, PPTX, 이미지, 웹 링크에서 본문을 추출한다."""

from __future__ import annotations

import io
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path

import requests
from bs4 import BeautifulSoup
from docx import Document
from PIL import Image

from src.extract_errors import ExtractError

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import trafilatura
except ImportError:  # pragma: no cover
    trafilatura = None

try:
    import pymupdf
except ImportError:  # pragma: no cover
    pymupdf = None


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".md",
    ".html",
    ".htm",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".tif",
    ".tiff",
    ".bmp",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}
USER_AGENT = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
    "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36"
)
SPARSE_PAGE_CHARS = 40


@dataclass
class ExtractedDocument:
    text: str
    source_name: str
    method: str
    file_bytes: bytes | None = None
    suffix: str = ""


def extract_from_bytes(filename: str, data: bytes) -> str:
    return load_file(filename, data).text


def extract_from_url(url: str, timeout: int = 25) -> str:
    return load_url(url, timeout=timeout).text


def load_file(
    filename: str,
    data: bytes,
    ocr_fallback: Callable | None = None,
) -> ExtractedDocument:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ExtractError(
            f"지원하지 않는 파일 형식입니다: {suffix or '(확장자 없음)'}. "
            "PDF, DOCX, PPTX, 이미지, TXT, MD, HTML만 사용할 수 있습니다."
        )
    if suffix == ".pdf":
        text, method = _from_pdf(data, ocr_fallback=ocr_fallback)
    elif suffix == ".docx":
        text, method = _from_docx(data), "docx"
    elif suffix == ".pptx":
        text, method = _from_pptx(data), "pptx"
    elif suffix in IMAGE_EXTENSIONS:
        from src.ocr import ocr_image_bytes

        text, method = ocr_image_bytes(data, fallback=ocr_fallback), "ocr"
    elif suffix in {".html", ".htm"}:
        text, method = _from_html(data.decode("utf-8", errors="replace")), "html"
    else:
        text, method = data.decode("utf-8", errors="replace"), "text"

    text = (text or "").strip()
    if not text:
        raise ExtractError("문서에서 텍스트를 찾지 못했습니다.")
    return ExtractedDocument(
        text=text,
        source_name=filename,
        method=method,
        file_bytes=data,
        suffix=suffix,
    )


def load_url(url: str, timeout: int = 25, ocr_fallback: Callable | None = None) -> ExtractedDocument:
    url = (url or "").strip()
    if not url:
        raise ExtractError("링크를 입력해 주세요.")
    if not url.startswith(("http://", "https://")):
        url = "https://" + url

    try:
        response = requests.get(
            url,
            timeout=timeout,
            headers={"User-Agent": USER_AGENT},
            allow_redirects=True,
        )
        response.raise_for_status()
    except requests.RequestException as exc:
        raise ExtractError(f"링크에서 문서를 가져오지 못했습니다: {exc}") from exc

    content_type = (response.headers.get("Content-Type") or "").lower()
    filename = _filename_from_url(url, content_type)
    suffix = Path(filename).suffix.lower()

    if "pdf" in content_type or suffix == ".pdf":
        doc = load_file(filename or "page.pdf", response.content, ocr_fallback=ocr_fallback)
        doc.source_name = url
        return doc
    if "word" in content_type or suffix == ".docx":
        doc = load_file(filename or "page.docx", response.content)
        doc.source_name = url
        return doc
    if content_type.startswith("image/") or suffix in IMAGE_EXTENSIONS:
        doc = load_file(filename or "page.png", response.content, ocr_fallback=ocr_fallback)
        doc.source_name = url
        return doc

    html = response.text
    text = ""
    if trafilatura is not None:
        text = trafilatura.extract(
            html,
            include_comments=False,
            include_tables=True,
            favor_recall=True,
        ) or ""
    if not text.strip():
        text = _from_html(html)
    if not text.strip():
        raise ExtractError("페이지에서 본문을 찾지 못했습니다.")
    return ExtractedDocument(
        text=text.strip(),
        source_name=url,
        method="url",
        file_bytes=None,
        suffix=".html",
    )


def load_pasted(text: str) -> ExtractedDocument:
    body = (text or "").strip()
    if not body:
        raise ExtractError("번역할 텍스트를 붙여 넣어 주세요.")
    return ExtractedDocument(
        text=body,
        source_name="붙여넣은 텍스트",
        method="paste",
        suffix=".txt",
    )


def _filename_from_url(url: str, content_type: str) -> str:
    path = url.split("?", 1)[0]
    name = Path(path).name
    if name:
        return name
    if "pdf" in content_type:
        return "download.pdf"
    if content_type.startswith("image/"):
        return "download.png"
    return "page.html"


def _from_pdf(data: bytes, ocr_fallback: Callable | None = None) -> tuple[str, str]:
    if pymupdf is not None:
        return _from_pdf_pymupdf(data, ocr_fallback=ocr_fallback)
    return _from_pdf_pypdf(data), "pdf"


def _from_pdf_pypdf(data: bytes) -> str:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractError(f"PDF를 읽지 못했습니다: {exc}") from exc
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    text = "\n\n".join(part for part in pages if part).strip()
    if not text:
        raise ExtractError("PDF에서 텍스트를 추출하지 못했습니다. 스캔본일 수 있습니다.")
    return text


def _from_pdf_pymupdf(data: bytes, ocr_fallback: Callable | None = None) -> tuple[str, str]:
    try:
        document = pymupdf.open(stream=data, filetype="pdf")
    except Exception as exc:  # noqa: BLE001
        raise ExtractError(f"PDF를 읽지 못했습니다: {exc}") from exc

    used_ocr = False
    pages: list[str] = []
    try:
        for index, page in enumerate(document, start=1):
            extracted = (page.get_text("text") or "").strip()
            if len(extracted) >= SPARSE_PAGE_CHARS:
                pages.append(extracted)
                continue
            ocr_text = _ocr_pdf_page(page, fallback=ocr_fallback)
            if ocr_text:
                used_ocr = True
                combined = "\n".join(part for part in (extracted, ocr_text) if part)
                pages.append(f"[페이지 {index}]\n{combined}")
            elif extracted:
                pages.append(extracted)
    finally:
        document.close()

    text = "\n\n".join(pages).strip()
    if not text:
        raise ExtractError("PDF에서 텍스트를 추출하지 못했습니다. 스캔본 OCR에도 실패했습니다.")
    return text, "pdf-ocr" if used_ocr else "pdf"


def _ocr_pdf_page(page, fallback: Callable | None = None) -> str:
    from src.ocr import ocr_pil

    pix = page.get_pixmap(matrix=pymupdf.Matrix(2, 2), alpha=False)
    image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    try:
        return ocr_pil(image, fallback=fallback)
    except ExtractError:
        return ""


def _from_docx(data: bytes) -> str:
    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractError(f"DOCX를 읽지 못했습니다: {exc}") from exc

    blocks: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            blocks.append(paragraph.text.strip())
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))
    text = "\n".join(blocks).strip()
    if not text:
        raise ExtractError("DOCX에서 텍스트를 찾지 못했습니다.")
    return text


def _from_pptx(data: bytes) -> str:
    if Presentation is None:
        raise ExtractError("PPTX 지원을 위해 python-pptx가 필요합니다.")
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ExtractError(f"PPTX를 읽지 못했습니다: {exc}") from exc

    blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            blocks.append(f"[슬라이드 {index}]\n" + "\n".join(texts))
    text = "\n\n".join(blocks).strip()
    if not text:
        raise ExtractError("PPTX에서 텍스트를 찾지 못했습니다.")
    return text


def _from_html(html: str) -> str:
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "form"]):
        tag.decompose()
    title = soup.title.get_text(" ", strip=True) if soup.title else ""
    body = soup.get_text("\n", strip=True)
    parts = [part for part in (title, body) if part]
    return "\n\n".join(parts).strip()
